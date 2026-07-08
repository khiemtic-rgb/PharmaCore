#!/usr/bin/env python3
"""Export sales-deck-v1.md -> PowerPoint in ấn (Novixa brand, screenshot + chú thích)."""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SCRIPT_DIR = Path(__file__).resolve().parent
MD_PATH = SCRIPT_DIR / "sales-deck-v1.md"
LOGO_PATH = SCRIPT_DIR.parents[2].joinpath("assets", "brand", "novixa-logo.png")
if not LOGO_PATH.exists():
    LOGO_PATH = Path(__file__).resolve().parents[4] / "client" / "admin" / "public" / "logo.png"

OUT_PATH = SCRIPT_DIR / "sales-deck-v1-with-screenshots.pptx"
OUT_PRINT = SCRIPT_DIR / "sales-deck-v1-print-brochure.pptx"
OUT_PATH_ALT = SCRIPT_DIR / "sales-deck-v1-full.pptx"

NAVY = RGBColor(0x1B, 0x3A, 0x6B)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK = RGBColor(0x1E, 0x29, 0x3B)
FOOTER = "Novixa · Hotline 0984.660.399 · khiemtic@gmail.com · KĐT Hồ Xương Rồng, P. Phan Đình Phùng, Thái Nguyên · novixa.vn"


def strip_md(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", text)
    for ch in ("✅", "❌", "🧪", "📋", "🌐", "📧", "📞", "📍"):
        text = text.replace(ch, "")
    return text.strip()


def parse_slides(md: str) -> list[dict]:
    slides: list[dict] = []
    main_md = md.split("## Phụ lục A")[0]
    blocks = re.split(r"\n---\n", main_md)

    for block in blocks:
        m = re.search(r"^## Slide (\d+) — (.+)$", block, re.MULTILINE)
        if not m:
            continue
        num = int(m.group(1))
        subtitle = m.group(2).strip()
        body = block[m.end() :].strip()
        lines = [ln.rstrip() for ln in body.splitlines()]

        title = subtitle
        lead = ""
        content_lines: list[str] = []
        callouts: list[str] = []
        notes: list[str] = []
        image: dict | None = None
        in_code = False

        for ln in lines:
            if not ln.strip():
                continue
            if ln.startswith("```"):
                in_code = not in_code
                continue
            if in_code:
                content_lines.append(strip_md(ln))
                continue
            if ln.startswith(">>"):
                callouts.append(strip_md(ln[2:].strip()))
                continue
            img = re.match(r"^!\[(.+?)\]\((.+?)\)\s*$", ln.strip())
            if img:
                image = {"alt": img.group(1), "rel": img.group(2)}
                continue
            if re.match(r"^\|.+\|$", ln):
                if re.match(r"^\|[-:\s|]+\|$", ln):
                    continue
                cells = [c.strip() for c in ln.strip("|").split("|")]
                if len(cells) >= 2:
                    row = " | ".join(strip_md(c) for c in cells if c.strip())
                    content_lines.append(row)
                continue
            if ln.startswith("- [ ]"):
                content_lines.append("[ ] " + strip_md(ln[5:].strip()))
                continue
            if ln.startswith("- "):
                content_lines.append(strip_md(ln[2:]))
                continue
            if re.match(r"^\d+\.\s", ln):
                content_lines.append(strip_md(ln))
                continue
            if ln.startswith("*") and ln.endswith("*") and not ln.startswith("**"):
                notes.append(strip_md(ln))
                continue
            bold = re.match(r"^\*\*(.+)\*\*$", ln)
            if bold:
                text = strip_md(ln)
                if not lead and num != 1:
                    lead = text
                elif num == 1 and title == subtitle:
                    title = text
                else:
                    content_lines.append(text)
                continue
            if ln.startswith("#"):
                continue
            content_lines.append(strip_md(ln))

        slides.append(
            {
                "num": num,
                "title": title,
                "lead": lead,
                "bullets": [x for x in content_lines if x],
                "callouts": callouts,
                "notes": "\n".join(notes),
                "image": image,
            }
        )

    slides.sort(key=lambda s: s["num"])
    return slides


def add_header_footer(slide, title: str, prs: Presentation) -> None:
    w = prs.slide_width
    h = prs.slide_height
    bar = slide.shapes.add_shape(1, 0, 0, w, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = f"  {title[:100]}"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    foot = slide.shapes.add_shape(1, 0, h - Inches(0.38), w, Inches(0.38))
    foot.fill.solid()
    foot.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
    foot.line.fill.background()
    ftf = foot.text_frame
    ftf.text = FOOTER
    fp = ftf.paragraphs[0]
    fp.font.size = Pt(8)
    fp.font.color.rgb = GRAY
    fp.alignment = PP_ALIGN.CENTER
    ftf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_text_box(slide, left, top, width, height, lines: list[str], *, font_size=13, bold_first=False, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text[:900]
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
        if bold_first and i == 0:
            p.font.bold = True
    return box


def add_callout(slide, left, top, width, lines: list[str]) -> None:
    if not lines:
        return
    h = Inches(0.28 + 0.22 * len(lines))
    shape = slide.shapes.add_shape(1, left, top, width, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
    shape.line.color.rgb = GREEN
    tf = shape.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_top = Inches(0.06)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line[:500]
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0x05, 0x46, 0x1A)
        if i == 0:
            p.font.bold = True


def add_cover(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w, h = prs.slide_width, prs.slide_height
    bg = slide.shapes.add_shape(1, 0, 0, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    accent = slide.shapes.add_shape(1, 0, h - Inches(0.12), w, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(4.9), Inches(0.55), height=Inches(1.35))

    title = data["bullets"][0] if data["bullets"] else "Novixa — Nền tảng quản trị nhà thuốc"
    add_text_box(
        slide,
        Inches(0.8),
        Inches(2.15),
        Inches(11.7),
        Inches(1.2),
        [title],
        font_size=32,
        bold_first=True,
        color=WHITE,
    )
    sub = data["bullets"][1:6] if len(data["bullets"]) > 1 else [
        "Smart Pharmacy Solutions",
        "ERP + POS + Kho lô FEFO + CRM + App khách + Báo cáo",
        "Founding Early Access 2026",
    ]
    add_text_box(slide, Inches(1.2), Inches(3.45), Inches(10.8), Inches(2.2), sub, font_size=16, color=RGBColor(0xBF, 0xDB, 0xFE))

    contact = [
        "Hotline 0984.660.399  ·  khiemtic@gmail.com",
        "KĐT Hồ Xương Rồng, phường Phan Đình Phùng, tỉnh Thái Nguyên",
        "www.novixa.vn",
    ]
    add_text_box(slide, Inches(1.2), Inches(5.85), Inches(10.8), Inches(1.0), contact, font_size=12, color=RGBColor(0x86, 0xEF, 0xAC))


def add_content_slide(prs: Presentation, data: dict, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, f"Slide {data['num']}/{total} — {data['title']}", prs)

    has_img = bool(data.get("image"))
    content_w = Inches(6.15) if has_img else Inches(12.2)
    body_lines: list[str] = []
    if data["lead"]:
        body_lines.append(data["lead"])
    body_lines.extend(data["bullets"])

    max_lines = 7 if has_img else 14
    overflow = body_lines[max_lines:]
    visible = body_lines[:max_lines]

    add_text_box(slide, Inches(0.45), Inches(0.95), content_w, Inches(5.5), visible, font_size=13)

    if has_img:
        img_meta = data["image"]
        img_path = (SCRIPT_DIR / img_meta["rel"]).resolve()
        img_top = Inches(0.95)
        if data.get("callouts"):
            add_callout(slide, Inches(6.72), img_top, Inches(6.15), data["callouts"])
            img_top += Inches(0.35 + 0.22 * len(data["callouts"]))
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), Inches(6.72), img_top, width=Inches(6.15))
            cap = slide.shapes.add_textbox(Inches(6.72), img_top + Inches(4.05), Inches(6.15), Inches(0.45))
            cp = cap.text_frame.paragraphs[0]
            cp.text = img_meta.get("alt", "")[:120]
            cp.font.size = Pt(9)
            cp.font.italic = True
            cp.font.color.rgb = GRAY
            cp.alignment = PP_ALIGN.CENTER
        else:
            overflow.append(f"[Thiếu ảnh: {img_meta['rel']}]")
    elif data.get("callouts"):
        add_callout(slide, Inches(0.45), Inches(5.95), Inches(12.2), data["callouts"])

    notes_slide = slide.notes_slide
    nf = notes_slide.notes_text_frame
    note_body = data["notes"]
    if overflow:
        note_body = (note_body + "\n\nThêm:\n" + "\n".join(overflow)).strip()
    nf.text = f"Slide {data['num']}/{total}\n\n{note_body}".strip()


def build_pptx(slides: list[dict], out_paths: list[Path]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    total = max(s["num"] for s in slides) if slides else 0
    for data in slides:
        if data["num"] > 40:
            continue
        if data["num"] == 1:
            add_cover(prs, data)
        else:
            add_content_slide(prs, data, total)

    saved = None
    for path in out_paths:
        try:
            prs.save(path)
            saved = path
            break
        except PermissionError:
            continue
    if saved is None:
        raise PermissionError("Cannot write PPT — close open PowerPoint files and retry.")
    return saved


def main() -> int:
    if not MD_PATH.exists():
        print(f"Missing: {MD_PATH}", file=sys.stderr)
        return 1
    md = MD_PATH.read_text(encoding="utf-8")
    slides = parse_slides(md)
    if not slides:
        print("No slides parsed.", file=sys.stderr)
        return 1
    saved = build_pptx(slides, [OUT_PRINT, OUT_PATH, OUT_PATH_ALT])
    n = len([s for s in slides if s["num"] <= 40])
    print(f"Wrote {n} slides -> {saved}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
